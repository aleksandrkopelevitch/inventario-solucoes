#!/usr/bin/env python3
"""
Builds resources/cati/cati-template.pptx from a real approved CATI deck.

Run once when the committee publishes a new template — not at request time.

Two things happen here, and the second is the one that matters:

1. Every slide is dropped, so what remains is the master, the theme and the
   six layouts. Orphaned parts (the deck's own photos, its notes slides) fall
   out of the package automatically once nothing references them.

2. The COVER and CLOSING decorations are lifted off the source slides and
   pinned onto their layouts. In the source deck those shapes — the corner
   brackets, the Leo CONECTADOS logo, the tagline band — were drawn on the
   slides themselves, so a stripped template produces a bare green rectangle.
   Moving them to the layout is what lets the renderer stay dumb: it places
   text and nothing else, and the identity comes from the template.

    build_template.py --source DECK.pptx --out resources/cati/cati-template.pptx
"""

import argparse
import copy

from pptx import Presentation
from pptx.oxml.ns import qn

# Text on the source cover/closing is that deck's own content ("SKBridge",
# "Fabio Caldart") — the renderer writes its own, so only the decoration is
# lifted.
DECORATION_TAGS = {qn('p:pic'), qn('p:grpSp'), qn('p:cxnSp'), qn('p:sp')}

RELATIONSHIP_ATTRS = (qn('r:embed'), qn('r:link'), qn('r:id'))


def is_text_box(element):
    """A <p:sp> that carries text is content, not decoration."""
    if element.tag != qn('p:sp'):
        return False
    texts = element.findall('.//' + qn('a:t'))
    return any((t.text or '').strip() for t in texts)


def copy_shape(src_part, dst_part, element):
    """
    Deep-copies a shape into another part, rewriting its relationship ids.

    An `r:embed` is only meaningful relative to the part that declares it, so
    copying the XML alone yields a picture pointing at whatever that id happens
    to mean in the destination — usually nothing, sometimes the wrong image.
    """
    new = copy.deepcopy(element)

    for el in new.iter():
        for attr in RELATIONSHIP_ATTRS:
            rId = el.get(attr)
            if not rId or rId not in src_part.rels:
                continue

            rel = src_part.rels[rId]

            if rel.is_external:
                new_rId = dst_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            else:
                new_rId = dst_part.relate_to(rel.target_part, rel.reltype)

            el.set(attr, new_rId)

    return new


def decorate_layout(prs, layout_name, source_slide):
    layout = next((l for l in prs.slide_layouts if l.name == layout_name), None)
    if layout is None:
        raise SystemExit(f'layout {layout_name!r} não existe no deck de origem')

    tree = layout.shapes._spTree
    moved = 0

    for shape in source_slide.shapes:
        element = shape._element
        if element.tag not in DECORATION_TAGS or is_text_box(element):
            continue
        tree.append(copy_shape(source_slide.part, layout.part, element))
        moved += 1

    print(f'  {layout_name!r}: +{moved} shapes')


def drop_all_slides(prs):
    id_list = prs.slides._sldIdLst
    for slide_id in list(id_list):
        prs.part.drop_rel(slide_id.rId)
        id_list.remove(slide_id)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--source', required=True, help='an approved CATI deck')
    parser.add_argument('--out', required=True)
    parser.add_argument('--cover-slide', type=int, default=1, help='1-based')
    parser.add_argument('--closing-slide', type=int, default=-1, help='1-based; -1 = last')
    args = parser.parse_args()

    prs = Presentation(args.source)
    slides = list(prs.slides)

    cover = slides[args.cover_slide - 1]
    closing = slides[args.closing_slide if args.closing_slide > 0 else len(slides) - 1]

    print(f'origem: {len(slides)} slides')
    decorate_layout(prs, cover.slide_layout.name, cover)
    decorate_layout(prs, closing.slide_layout.name, closing)

    drop_all_slides(prs)
    prs.save(args.out)

    check = Presentation(args.out)
    print(f'template: {len(check.slides)} slides, {len(check.slide_layouts)} layouts')
    for layout in check.slide_layouts:
        print(f'  {layout.name!r}: {len(list(layout.shapes))} shapes')


if __name__ == '__main__':
    main()
