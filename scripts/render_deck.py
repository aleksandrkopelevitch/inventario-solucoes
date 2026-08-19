#!/usr/bin/env python3
"""
Renders a CATI deck spec into the corporate .pptx template.

Called by App\\Actions\\Cati\\RenderSubmissionDeck through Symfony Process.
It is deliberately dumb: it places what the spec says and decides nothing.
Everything about WHICH slides exist and what they say is decided in PHP
(BuildDeckSpec) and checked before it gets here (DeckSpecValidator).

Why a template and not a deck built from scratch: the layouts carry the Leo
identity — the title placeholder, the footer band, the dark green (#144227)
and the theme font. A deck that comes out only *almost* corporate is the
failure mode most likely to make someone redo it by hand.

    render_deck.py --template T.pptx --spec spec.json --out deck.pptx
"""

import argparse
import json
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# Layouts are resolved BY NAME, not by index: an index silently renders the
# wrong layout if the committee ever reorders the template, while a name that
# stops existing fails loudly. Mirrors the warning in resources/cati/README.md.
LAYOUTS = {
    'cover': 'Fundo escuto sólido',
    'content': 'Fundo branco com título',
    'closing': 'Fundo escuro limpo com objetos',
}

# Measured off a real approved deck (CATI_SKBridge_SkyMob), so generated slides
# land exactly where hand-made ones do. The content layout has ONLY a title
# placeholder — there is no body placeholder to inherit from, which is why the
# real decks all carry a hand-drawn text box at this position.
BODY = dict(left=Inches(0.5), top=Inches(0.96), width=Inches(12.33), height=Inches(5.91))

# Body type scale, also measured: 16pt for a top-level line, stepping down.
SIZES = {0: Pt(16), 1: Pt(14), 2: Pt(12), 3: Pt(11)}

# The dark layouts have no placeholders, so their text boxes are ours — and a
# plain text box inherits the master's dark body colour, which is invisible on
# #144227. Cover and closing set it explicitly; content slides inherit.
LIGHT_TEXT = RGBColor(0xFF, 0xFF, 0xFF)

BULLET = '•'

# Measured off the approved deck: a bulleted line indents 285750 EMU (0.3125in)
# with a matching negative first-line indent, so the text hangs level under
# itself. A line WITHOUT a bullet says so explicitly (`a:buNone`) — inheriting
# the master's list style would otherwise put a glyph in front of a paragraph.
BULLET_INDENT = 285750

# The real decks set 10pt in table cells, header bold.
TABLE_HEADER_SIZE = Pt(10)
TABLE_BODY_SIZE = Pt(10)


def layout_by_name(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise SystemExit(
        f"layout {name!r} não existe no template (encontrados: "
        f"{[l.name for l in prs.slide_layouts]})"
    )


def set_title(slide, text):
    """Fills the layout's title placeholder, which is where the corporate 24pt lives."""
    if slide.shapes.title is None:
        return False
    slide.shapes.title.text = text
    return True


def add_text_box(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def set_bullet(paragraph, level, bulleted):
    """
    Real list formatting rather than a `•` typed into the text.

    Element order inside <a:pPr> is fixed by the schema — a bullet appended
    after <a:defRPr> is invalid and PowerPoint offers to repair the file —
    so this inserts before whatever may follow it rather than appending.
    """
    pPr = paragraph._p.get_or_add_pPr()

    for tag in ('a:buNone', 'a:buChar', 'a:buAutoNum'):
        for existing in pPr.findall(qn(tag)):
            pPr.remove(existing)

    if bulleted:
        pPr.set('marL', str(BULLET_INDENT * (level + 1)))
        pPr.set('indent', str(-BULLET_INDENT))
        element = pPr.makeelement(qn('a:buChar'), {'char': BULLET})
    else:
        pPr.set('marL', str(BULLET_INDENT * level))
        pPr.set('indent', '0')
        element = pPr.makeelement(qn('a:buNone'), {})

    pPr.insert_element_before(element, 'a:tabLst', 'a:defRPr', 'a:extLst')


def write_blocks(frame, blocks, dark=False):
    """Bullets and paragraphs into one text frame; tables are handled by the caller."""
    first = True

    for block in blocks:
        if block.get('type') == 'table':
            continue

        para = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False

        level = int(block.get('level', 0) or 0)
        para.level = min(level, 4)

        run = para.add_run()
        run.text = block.get('text', '')
        run.font.size = SIZES.get(level, SIZES[3])
        if dark:
            run.font.color.rgb = LIGHT_TEXT

        para.space_after = Pt(6)
        set_bullet(para, level, block.get('type') == 'bullet')


def add_table(slide, block, left, top, width, height):
    columns = block['columns']
    rows = block['rows']
    shape = slide.shapes.add_table(len(rows) + 1, len(columns), left, top, width, height)
    table = shape.table

    for c, label in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = str(label)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = TABLE_HEADER_SIZE
                run.font.bold = True

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = TABLE_BODY_SIZE

    return shape


def add_image(slide, block, left, top, width, height):
    """
    A diagram, rendered by the F3 canvas and published on every layout save.

    It goes in as a picture with a hyperlink back to the canvas on purpose:
    native shapes would make the deck a SECOND place the diagram can be
    edited, and the two would drift apart on the first meeting.
    """
    picture = slide.shapes.add_picture(block['path'], left, top, height=height)

    # add_picture with only a height keeps the aspect ratio; centre whatever
    # width that produced inside the body box.
    if picture.width < width:
        picture.left = left + int((width - picture.width) / 2)
    else:
        # Too wide: re-fit by width instead, and re-centre vertically.
        scale = width / picture.width
        picture.width = int(width)
        picture.height = int(picture.height * scale)
        picture.left = left
        picture.top = top + int((height - picture.height) / 2)

    if block.get('link'):
        picture.click_action.hyperlink.address = block['link']

    return picture


def render_cover(slide, spec_slide):
    frame = add_text_box(slide, Inches(0.41), Inches(3.24), Inches(8.73), Inches(0.9))
    run = frame.paragraphs[0].add_run()
    run.text = spec_slide['title']
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = LIGHT_TEXT

    subtitle = ' '.join(filter(None, [spec_slide.get('subtitle'), spec_slide.get('footnote')]))
    if subtitle:
        frame2 = add_text_box(slide, Inches(0.41), Inches(4.14), Inches(8.73), Inches(0.68))
        run2 = frame2.paragraphs[0].add_run()
        run2.text = subtitle
        run2.font.size = Pt(20)
        run2.font.color.rgb = LIGHT_TEXT


def render_closing(slide, spec_slide):
    frame = add_text_box(slide, Inches(3.17), Inches(4.23), Inches(7.0), Inches(0.8))
    run = frame.paragraphs[0].add_run()
    run.text = spec_slide['title']
    run.font.size = Pt(44)
    run.font.color.rgb = LIGHT_TEXT

    if spec_slide.get('subtitle'):
        frame2 = add_text_box(slide, Inches(3.17), Inches(5.05), Inches(7.0), Inches(0.68))
        run2 = frame2.paragraphs[0].add_run()
        run2.text = spec_slide['subtitle']
        run2.font.size = Pt(20)
        run2.font.color.rgb = LIGHT_TEXT


def render_content(slide, spec_slide):
    set_title(slide, spec_slide['title'])

    blocks = spec_slide.get('blocks', [])
    text_blocks = [b for b in blocks if b.get('type') in ('bullet', 'paragraph')]
    figures = [b for b in blocks if b.get('type') in ('table', 'image')]

    top = BODY['top']
    height = BODY['height']

    if text_blocks:
        # Text first, then whatever figure follows it — the same split the real
        # decks use on a table slide (intro line above, table below).
        text_height = Inches(0.61) if figures else height
        frame = add_text_box(slide, BODY['left'], top, BODY['width'], text_height)
        write_blocks(frame, text_blocks)
        top = top + text_height + Inches(0.14)
        height = height - text_height - Inches(0.14)

    # One figure per slide, guaranteed upstream by BuildDeckSpec::paginate().
    # Anything beyond the first would have to overlap it, so this refuses
    # rather than silently dropping it — a spec that gets here with two
    # figures on one slide is a builder bug, not a rendering choice.
    if len(figures) > 1:
        raise SystemExit(
            f"slide {spec_slide['title']!r} traz {len(figures)} figuras; "
            'só cabe uma por slide (ver BuildDeckSpec::paginate).'
        )

    for figure in figures:
        if figure['type'] == 'table':
            add_table(slide, figure, BODY['left'], top, BODY['width'], height)
        else:
            add_image(slide, figure, BODY['left'], top, BODY['width'], height)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--template', required=True)
    parser.add_argument('--spec', required=True)
    parser.add_argument('--out', required=True)
    args = parser.parse_args()

    with open(args.spec, encoding='utf-8') as handle:
        spec = json.load(handle)

    prs = Presentation(args.template)

    for spec_slide in spec.get('slides', []):
        kind = spec_slide.get('layout')
        if kind not in LAYOUTS:
            raise SystemExit(f'layout desconhecido no spec: {kind!r}')

        slide = prs.slides.add_slide(layout_by_name(prs, LAYOUTS[kind]))

        if kind == 'cover':
            render_cover(slide, spec_slide)
        elif kind == 'closing':
            render_closing(slide, spec_slide)
        else:
            render_content(slide, spec_slide)

    prs.save(args.out)
    print(json.dumps({'slides': len(spec.get('slides', []))}))


if __name__ == '__main__':
    sys.exit(main())
